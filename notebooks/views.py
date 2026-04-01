# Python utitlities
import os
import re
import uuid
import time
import random
import shutil
import stat
import logging
from threading import Thread
from datetime import datetime
from pathlib import Path

import requests as http_requests
import secrets
import urllib.parse

# Django web framework tools
from django.conf import settings
from django.http import JsonResponse, HttpResponseNotAllowed
from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth import authenticate, login, logout, get_user_model
from django.contrib.auth.forms import PasswordResetForm, SetPasswordForm
from django.contrib.auth.tokens import default_token_generator
from django.contrib.auth.password_validation import validate_password
from django.core.exceptions import ValidationError
from django.utils.http import urlsafe_base64_decode
from django.utils.encoding import force_str
from django.views.decorators.http import require_GET, require_POST
from django.views.decorators.cache import never_cache
from django.core.cache import cache

# Document parsing
from pypdf import PdfReader
import docx
from pptx import Presentation

# AI pipeline
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import CharacterTextSplitter
from openai import OpenAI

# Email sending
from django.core.mail import send_mail

# Logging
log = logging.getLogger("timing")
log.setLevel(logging.INFO)
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s"
)
User = get_user_model()


# Helpers


def allowed_file(filename: str) -> bool:
    allowed = {"pdf", "docx", "txt", "pptx"}
    return "." in filename and filename.rsplit(".", 1)[1].lower() in allowed


def secure_filename(name: str) -> str:
    name = os.path.basename(name)
    name = re.sub(r"[^A-Za-z0-9.\-_]+", "_", name)
    return name.strip("._") or f"upload_{uuid.uuid4().hex}"


def extract_text_from_pdf(file_path: str) -> str:
    try:
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            parts = []
            page_num = 1
            for page in reader.pages:
                text = page.extract_text() or ""
                lines = [line.strip() for line in text.splitlines() if line.strip()]

                parts.append(f"\n--- Page {page_num} ---\n")

                for line in lines:
                    if len(line) < 80 and line.isupper() and not line.endswith("."):
                        parts.append(f"\n--- {line} ---\n")
                    else:
                        parts.append(line)
                page_num += 1
            return "\n".join(parts)
    except Exception as e:
        return f"Error extracting text from PDF: {e}"


def extract_text_from_docx(file_path: str) -> str:
    try:
        d = docx.Document(file_path)
        parts = []
        for p in d.paragraphs:
            if p.style and p.style.name.startswith("Heading"):
                parts.append(f"\n --- {p.text.upper()} ---\n")
            elif p.text.strip():
                parts.append(p.text)
        return "\n".join(parts)
    except Exception as e:
        return f"Error extracting text from DOCX: {e}"


def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"\n--- Slide {i} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text.strip())
        return "\n".join(parts).strip()
    except Exception as e:
        return f"Error extracting text from PPTX: {e}"


def process_uploaded_file(django_file):
    filename = secure_filename(django_file.name)
    save_path = Path(settings.UPLOAD_FOLDER) / filename

    # Save uploaded file
    with open(save_path, "wb") as out:
        for chunk in django_file.chunks():
            out.write(chunk)

    lower = filename.lower()
    if lower.endswith(".pdf"):
        text = extract_text_from_pdf(str(save_path))
    elif lower.endswith(".docx"):
        text = extract_text_from_docx(str(save_path))
    elif lower.endswith(".pptx"):
        text = extract_text_from_pptx(str(save_path))
    elif lower.endswith(".txt"):
        text = save_path.read_text(encoding="utf-8", errors="ignore")
    else:
        text = "Unsupported file type."

    base_name_without_ext = os.path.splitext(filename)[0]
    return text, base_name_without_ext, filename


def get_openai_client():
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY not set in environment.")
    return OpenAI(api_key=api_key)


def get_document_prompt(docs):
    out = []
    for i, d in enumerate(docs, 1):
        text = d if isinstance(d, str) else getattr(d, "page_content", "")
        out.append(f"\nContent {i}:\n{text}\n")
    return "\n".join(out)


def _on_rm_error(func, path, exc_info):
    try:
        os.chmod(path, stat.S_IWRITE)
        func(path)
    except Exception:
        logging.warning("Could not remove: %s", path)


def _get_safe_redirect_url(request):
    next_url = (request.POST.get("next") or request.GET.get("next") or "").strip()
    if next_url.startswith("/"):
        return next_url
    return request.META.get("HTTP_REFERER") or "/"


def _get_client_ip(request):
    forwarded = request.META.get("HTTP_X_FORWARDED_FOR")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.META.get("REMOTE_ADDR", "unknown")


def _is_rate_limited(request, action, max_attempts):
    """Read-only check — returns True if this IP is over the limit."""
    ip = _get_client_ip(request)
    return cache.get(f"rl:{action}:{ip}", 0) >= max_attempts


def _record_attempt(request, action, window_seconds):
    """
    Increment the attempt counter for this IP.
    Uses cache.add on the first attempt so the window is fixed from that
    point. Subsequent calls use cache.incr which does not reset the timeout.
    """
    ip = _get_client_ip(request)
    key = f"rl:{action}:{ip}"
    if not cache.add(key, 1, timeout=window_seconds):
        try:
            cache.incr(key)
        except ValueError:
            pass  # key expired between add and incr — harmless


def _process_job(
    job_id: str,
    text: str,
    persist_dir: str,
    filename: str,
    start_pct: int = 10,
    end_pct: int = 100,
):

    _update_floor = [0]

    def update(phase, pct, **extra):
        safe_pct = max(pct, _update_floor[0])
        _update_floor[0] = safe_pct
        cache.set(
            f"job:{job_id}", {"phase": phase, "pct": safe_pct, **extra}, timeout=3600
        )

    try:
        update("Preparing", 2)
        os.makedirs(persist_dir, exist_ok=True)

        # Split text - 2% -> 10%
        update("Processing", 5)
        text_splitter = CharacterTextSplitter(
            separator=" ", chunk_size=5000, chunk_overlap=100
        )
        docs = text_splitter.split_text(text) if text else []
        update("processing", 10)

        # Build embeddings + DB - 10% -> 15%
        update("Processing", 12)
        client = get_openai_client()
        embeddings = OpenAIEmbeddings(
            model="text-embedding-3-large", openai_api_key=client.api_key
        )
        vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)
        update("processing", 15)

        # Embed chunks - 15% -> 75%
        total = max(len(docs), 1)
        batch = 25
        added = 0
        for i in range(0, len(docs), batch):
            chunk = docs[i : i + batch]

            # Tick upward BEFORE the API call so UI feels responsive
            pre_pct = 15 + int(60 * added / total)
            update("processing", pre_pct)

            vectordb.add_texts(chunk)
            added += len(chunk)

            # Tick again AFTER so the bar moves on completion too
            post_pct = 15 + int(60 * added / total)
            update("processing", post_pct)

        # Retrieving content — 75% → 80%
        update("processing", 77)
        raw = vectordb.get(include=["documents"])
        sample = (raw.get("documents") or [])[:15]
        prompt = get_document_prompt(sample) if sample else "No content available."
        update("processing", 80)

        # Summarizing — 80% → 95% (streamed ticks while waiting)
        update("processing", 82)

        system_message = (
            f"Generate a summary of the following notebook content::\n\n"
            f"\n\n###\n{prompt}\n###\n\n"
            "The summary should contain the title of the book and a short sentence about the notebook"
            "The first line must be the notebook title, wrapped in double asterisks, like:\n"
            "Title: at the beginning of the notebook title"
            "Then add two newline characters (\\n\\n).\n"
            "After that, write the rest of the summary"
            "The summary should never be more that 2 sentences"
            "Be precise, avoid opinions, and summarize the main points in a clear and structured way. "
            "If the document has multiple sections, break it into meaningful segments."
        )

        # Use streaming so we can tick progress while waiting for OpenAI
        update("Summarizing", 85)
        summary_text = ""
        pct = 85
        with get_openai_client().chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": system_message}],
            temperature=0.2,
            stream=True,  # ← stream tokens
        ) as stream:
            for chunk in stream:
                delta = chunk.choices[0].delta.content or ""
                summary_text += delta
                if pct < 95:
                    pct += 1  # increment 1% per token batch — smooth crawl to 95
                    update("Summarizing", pct)

        # Done
        update("completed", 100, summary=summary_text, filename=filename)

    except Exception as e:
        cache.set(
            f"job:{job_id}", {"phase": "error", "pct": 0, "error": str(e)}, timeout=3600
        )


# ---------- Views (routes) ----------


@require_GET
def home(request):
    return render(request, "home.html")


@never_cache
@require_POST
def register_user(request):
    redirect_url = _get_safe_redirect_url(request)

    if _is_rate_limited(request, "register", max_attempts=5):
        messages.error(request, "Too many sign-up attempts. Please try again in an hour.", extra_tags="auth-signup")
        return redirect(redirect_url)

    _record_attempt(request, "register", window_seconds=3600)

    name = (request.POST.get("name") or "").strip()
    email = (request.POST.get("email") or "").strip().lower()
    password = request.POST.get("password") or ""

    if not name or not email or not password:
        messages.error(request, "Enter name, email, and password.", extra_tags="auth-signup")
        return redirect(redirect_url)

    if User.objects.filter(email__iexact=email).exists():
        messages.error(request, "Email already registered.", extra_tags="auth-signup")
        return redirect(redirect_url)

    first_name, _, last_name = name.partition(" ")
    user = User(username=email, email=email, first_name=first_name, last_name=last_name.strip())

    try:
        validate_password(password, user=user)
    except ValidationError as exc:
        error_text = " ".join(exc.messages)
        if "too short" in error_text.lower():
            error_text = "Password must be at least 8 characters."
        elif "too common" in error_text.lower():
            error_text = "Choose a less common password."
        elif "entirely numeric" in error_text.lower():
            error_text = "Password cannot be only numbers."
        messages.error(request, error_text, extra_tags="auth-signup")
        return redirect(redirect_url)

    user.set_password(password)
    user.save()
    messages.success(
        request,
        "Account created. Please sign in.",
        extra_tags="auth-login auth-signup-success",
    )
    return redirect(redirect_url)


@never_cache
@require_POST
def login_user(request):
    redirect_url = _get_safe_redirect_url(request)

    if _is_rate_limited(request, "login", max_attempts=10):
        messages.error(request, "Too many login attempts. Please try again in 15 minutes.", extra_tags="auth-login")
        return redirect(redirect_url)

    email = (request.POST.get("email") or "").strip().lower()
    password = request.POST.get("password") or ""
    remember = request.POST.get("remember")

    if not email or not password:
        messages.error(request, "Enter email and password.", extra_tags="auth-login")
        return redirect(redirect_url)

    user = authenticate(request, username=email, password=password)
    if user is None:
        _record_attempt(request, "login", window_seconds=900)  # only count failures
        messages.error(request, "Invalid email or password.", extra_tags="auth-login")
        return redirect(redirect_url)

    login(request, user)
    if not remember:
        request.session.set_expiry(0)
    return redirect(redirect_url)


@never_cache
@require_POST
def logout_user(request):
    logout(request)
    messages.success(request, "You have been signed out.")
    return redirect(_get_safe_redirect_url(request))


def _get_upload_page_context(request):
    job_id = request.session.get("job_id")
    if job_id:
        st = cache.get(f"job:{job_id}") or {}
        phase = (st.get("phase") or "").lower()

        if phase == "completed" and st.get("summary"):
            docs = request.session.get("docs", {})
            filename = st.get("filename") or request.session.get("uploaded_filename")

            if filename:
                info = docs.get(filename, {})
                info["persist_dir"] = info.get("persist_dir") or request.session.get(
                    "persist_directory"
                )
                info["summary"] = st["summary"]
                docs[filename] = info
                request.session["docs"] = docs

                request.session["uploaded_filename"] = filename
                request.session["summary_text"] = st["summary"]

            cache.delete(f"job:{job_id}")
            request.session.pop("job_id", None)
            job_id = None

        elif phase == "error":
            messages.error(
                request,
                f"⚠️ Could not build index or generate summary: {st.get('error')}",
            )
            cache.delete(f"job:{job_id}")
            request.session.pop("job_id", None)
            job_id = None

    return {
        "filename": request.session.get("uploaded_filename"),
        "summary": request.session.get("summary_text"),
        "job_id": job_id,
    }


def privacy_policy(request):
    return render(request, "privacy_policy.html")


def terms_of_service(request):
    return render(request, "terms_of_service.html")


@require_POST
def delete_doc(request):
    try:
        data = getattr(request, "json", None)
    except Exception:
        data = None

    # Django doesn’t auto-parse JSON. Do it safely:
    import json
    payload = {}
    if request.body:
        try:
            payload = json.loads(request.body.decode("utf-8"))
        except Exception:
            payload = {}

    filename = (payload.get("filename") or request.POST.get("filename") or "").strip()
    if not filename:
        return JsonResponse({"ok": False, "error": "Please select a document to delete."}, status=400)

    docs = request.session.get("docs", {})
    if filename not in docs:
        return JsonResponse({"ok": False, "error": "Selected document not found."}, status=404)

    persist_dir = docs[filename].get("persist_dir")
    if not persist_dir:
        return JsonResponse({"ok": False, "error": "No database path stored for this document."}, status=500)

    persist_dir = os.path.abspath(persist_dir)
    if os.path.isdir(persist_dir):
        last_err = None
        for attempt in range(3):
            try:
                shutil.rmtree(persist_dir, onerror=_on_rm_error)
                last_err = None
                break
            except Exception as e:
                last_err = e
                time.sleep(0.3)

        if last_err:
            import logging
            logging.getLogger(__name__).error("Failed to delete database at %s: %s", persist_dir, last_err)
            return JsonResponse({"ok": False, "error": "Failed to delete the database. Please try again."}, status=500)

    docs.pop(filename, None)
    request.session["docs"] = docs

    if request.session.get("uploaded_filename") == filename:
        request.session.pop("uploaded_filename", None)
        request.session.pop("summary_text", None)
        request.session.pop("persist_directory", None)

    # Also remove the uploaded source file from disk and session list — safely
    try:
        uploaded_files = request.session.get("uploaded_files", [])
        if filename in uploaded_files:
            uploaded_files.remove(filename)
            request.session["uploaded_files"] = uploaded_files

        upload_base = Path(settings.UPLOAD_FOLDER).resolve()
        upload_path = (upload_base / filename).resolve()

        # Ensure we never delete files outside the configured upload folder.
        try:
            inside = upload_path.is_relative_to(upload_base)
        except AttributeError:
            # Fallback for older Python: compare common prefixes
            inside = str(upload_path).startswith(str(upload_base) + os.sep) or str(upload_path) == str(upload_base)

        if inside and upload_path.exists():
            upload_path.unlink()
        else:
            logging.getLogger(__name__).warning(
                "Refused to delete upload outside upload folder: %s", upload_path
            )
    except Exception as e:
        logging.getLogger(__name__).exception("Failed to delete uploaded file %s: %s", filename, e)

    messages.success(request, f"Deleted '{filename}' successfully.")
    return JsonResponse({"ok": True, "message": f"Deleted '{filename}' successfully."})

@require_GET
def upload_notebook(request):
    return render(request, "upload_notebook.html", _get_upload_page_context(request))


@require_GET
def get_summary(request):
    filename = request.GET.get("filename", "")
    docs = request.session.get("docs", {})
    info = docs.get(filename) or {}
    summary = info.get("summary") or ""
    return JsonResponse({"ok": True, "summary": summary})


@require_POST
def init_upload(request):
    job_id = uuid.uuid4().hex
    cache.set(f"job:{job_id}", {"phase": "Uploading", "pct": 1}, timeout=3600)
    request.session["job_id"] = job_id
    return JsonResponse({"ok": True, "job_id": job_id})


@require_POST
def upload(request):
    job_id = (
        request.headers.get("X-Job-Id")
        or request.session.get("job_id")
        or uuid.uuid4().hex
    )
    request.session["job_id"] = job_id
    cache.set(f"job:{job_id}", {"phase": "Uploading", "pct": 1}, timeout=3600)

    is_xhr = request.headers.get("X-Requested-With") == "XMLHttpRequest"

    if "file" not in request.FILES:
        msg = "No file part in request."
        if is_xhr:
            return JsonResponse({"ok": False, "error": msg}, status=400)
        messages.error(request, msg)
        return redirect("upload_notebook")

    f = request.FILES["file"]
    if not f.name:
        msg = "No file selected."
        if is_xhr:
            return JsonResponse({"ok": False, "error": msg}, status=400)
        messages.error(request, msg)
        return redirect("upload_notebook")

    if not allowed_file(f.name):
        msg = "Unsupported file type. Please upload PDF, DOCX, PPTX or TXT."
        if is_xhr:
            return JsonResponse({"ok": False, "error": msg}, status=400)
        messages.error(request, msg)
        return redirect("upload_notebook")

    text, base, filename = process_uploaded_file(f)

    cache.set(f"job:{job_id}", {"phase": "queued", "pct": 40, "filename": filename})

    uploaded = request.session.get("uploaded_files", [])
    uploaded.append(filename)
    request.session["uploaded_files"] = uploaded

    persist_dir = os.path.abspath(f"./chroma_db_{base}_{uuid.uuid4().hex[:8]}")

    docs = request.session.get("docs", {})
    docs[filename] = {"persist_dir": persist_dir}
    request.session["docs"] = docs
    request.session["persist_directory"] = persist_dir
    request.session["uploaded_filename"] = filename
    request.session["summary_text"] = None
    request.session["summary_generated"] = False

    t = Thread(
        target=_process_job,
        args=(job_id, text, persist_dir, filename, 10, 100),
        daemon=True,
    )
    t.start()

    if is_xhr:
        return JsonResponse(
            {"ok": True, "job_id": job_id, "filename": filename}, status=200
        )
    return redirect("upload_notebook")


@require_GET
def get_progress(request, job_id):
    st = cache.get(f"job:{job_id}")
    if not st:
        return JsonResponse(
            {"ok": False, "missing": True, "phase": "missing", "pct": 0}, status=404
        )
    resp = JsonResponse({"ok": True, **st})
    resp["Cache-Control"] = "no-store"
    return resp


@require_POST
def ask(request):
    import json

    payload = {}
    if request.body:
        try:
            payload = json.loads(request.body.decode("utf-8"))
        except Exception:
            payload = {}

    question = (payload.get("question") or "").strip()
    filename = payload.get("filename")

    if not question:
        return JsonResponse({"ok": False, "error": "Question is required."}, status=400)

    docs = request.session.get("docs", {})
    info = docs.get(filename or "", {})
    persist_dir = info.get("persist_dir") if info else None

    if not persist_dir or not os.path.isdir(persist_dir):
        return JsonResponse(
            {
                "ok": False,
                "error": "Please select a Notebook before asking a Question.",
            },
            status=400,
        )

    client = get_openai_client()
    embeddings = OpenAIEmbeddings(
        model="text-embedding-3-large", openai_api_key=client.api_key
    )
    vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)

    retrieved = vectordb.similarity_search(question, k=10)
    context = get_document_prompt(retrieved)

    system_message = (
        "You are a professor teaching a course. Use the following notebook content "
        f"to answer student questions accurately and concisely:\n\n{context}\n\n"
        "Be precise and avoid opinions."
        "Only state what is in the notebook content"
        "Do not state what is not in the given notebook and be very precise and straight forward "
    )

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_message},
            {"role": "user", "content": question},
        ],
        temperature=0.1,
    )
    answer = resp.choices[0].message.content
    return JsonResponse({"ok": True, "answer": answer})


class Timer:
    def __init__(self, name):
        self.name = name
        self.start = time.perf_counter()

    def done(self, extra=""):
        elapsed = time.perf_counter() - self.start
        log.info("[TIMER] %-20s %7.3fs %s", self.name, elapsed, extra)
        return elapsed


@require_POST
def generate_quiz(request):
    t_total = Timer("generate_quiz TOTAL")

    import json

    payload = {}
    if request.body:
        try:
            payload = json.loads(request.body.decode("utf-8"))
        except Exception:
            payload = {}

    t_request = Timer("parse request")
    num = int(payload.get("num_questions", 5))
    filename = payload.get("filename")
    t_request.done(f"(num={num})")

    t_session = Timer("session lookup")
    docs = request.session.get("docs", {})
    info = docs.get(filename or "", {})
    persist_dir = info.get("persist_dir") if info else None
    t_session.done()

    if not persist_dir or not os.path.isdir(persist_dir):
        return JsonResponse(
            {"ok": False, "error": "Please select a Notebook before generating Quiz."},
            status=400,
        )

    t_vectordb = Timer("load vector DB")
    client = get_openai_client()
    embeddings = OpenAIEmbeddings(
        model="text-embedding-3-large", openai_api_key=client.api_key
    )
    vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)
    t_vectordb.done()

    t_fetch = Timer("fetch documents")
    raw = vectordb.get(include=["documents"])
    all_docs = raw.get("documents", [])
    t_fetch.done(f"(docs={len(all_docs)})")

    t_sample = Timer("sample documents")
    num_samples = min(20, len(all_docs)) if all_docs else 0
    sample = random.sample(all_docs, num_samples) if num_samples > 0 else []
    context = get_document_prompt(sample) if sample else "No content available."
    t_sample.done(f"(sampled={num_samples}, chars={len(context)})")

    t_prompt = Timer("build prompt")
    system_message = (
        f"Generate {num} multiple-choice quiz questions from the following notebook content: "
        f"\n\n###\n{context}\n###\n\n"
        "Provide the questions in the same language as the notebook content. Make sure the questions are clear, concise, and directly related to the content. "
        "Each question should have 5 answer choices (A,B,C,D,E) and indicate the correct answer at the end:"
        "IMPORTANT: Distribute the correct answers evenly and unpredictably across A, B, C, D, and E. "
        "Do NOT default to B as the correct answer. Ensure that across all questions, each letter (A, B, C, D, E) "
        "appears as the correct answer roughly equally. Vary the position of the correct answer deliberately. "
        f"For {num} questions, spread correct answers so no single letter dominates. "
        """The format of the reply should be strictly:
        Question 1: <question>
        A)  <answer choice A>
        B)  <answer choice B>
        C)  <answer choice C>
        D)  <answer choice D>
        E)  <answer choice E>
        Correct Answer: <letter only>
        Explanation: <one sentence explaining why the correct answer is right>

        Question 2: <question>
         ..."""
    )
    t_prompt.done()

    t_llm = Timer("LLM generation")
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": system_message}],
        temperature=0.2,
    )
    t_llm.done()

    t_parse = Timer("parse LLM output")
    text = resp.choices[0].message.content.strip()
    print("RAW LLM OUTPUT:\n", text)
    blocks = re.split(r"\n(?=Question\s*\d+)", text.strip())
    quiz = []
    for b in blocks:
        lines = [l for l in b.splitlines() if l.strip()]
        if len(lines) >= 7 and lines[0].lower().startswith("question"):
            q = lines[0].strip()
            choices = lines[1:6]
            correct_letter = lines[6].split(":")[-1].strip()

            explanation = ""
            if len(lines) >= 8 and lines[7].lower().startswith("explanation"):
                explanation = lines[7].split(":", 1)[-1].strip()

            # Mapping letter to the actual correct answer text
            letter_map = {chr(65 + i): choice for i, choice in enumerate(choices)}
            correct_text = letter_map.get(correct_letter, choices[0])

            # Shuffle the choices
            random.shuffle(choices)

            # Find where the correct answer landed after shuffle
            new_correct_letter = chr(
                65 + choices.index(correct_text)
            )  # A, B, C, D, or E

            # Strip old letter prefixes and re-label A-E
            relabeled = [
                f"{chr(65+i)}) {c.split(')', 1)[-1].strip()}"
                for i, c in enumerate(choices)
            ]

            quiz.append(
                {"question": q, "choices": relabeled, "correct": new_correct_letter, "explanation": explanation}
            )
    t_parse.done(f"(parsed={len(quiz)})")

    t_total.done()
    # Do not expose explanations to anonymous users — strip them server-side
    if not request.user.is_authenticated:
        for q in quiz:
            q["explanation"] = ""

    return JsonResponse({"ok": True, "quiz": quiz})


@require_POST
def save_result(request):
    import json

    payload = {}
    if request.body:
        try:
            payload = json.loads(request.body.decode("utf-8"))
        except Exception:
            payload = {}

    filename = payload.get("filename")
    correct = payload.get("correct")
    total = payload.get("total")
    percent = payload.get("percent")

    if not filename or correct is None or total is None:
        return JsonResponse({"ok": False, "error": "Missing result data"}, status=400)

    results = request.session.get("results", [])
    correct_i = int(correct)
    total_i = int(total)
    percent_i = (
        int(percent) if percent is not None else round((correct_i / total_i) * 100)
    )

    results.insert(
        0,
        {
            "filename": filename,
            "correct": correct_i,
            "total": total_i,
            "percent": percent_i,
            # store both a human-friendly server-local string and an ISO UTC timestamp
            "test_datetime": datetime.now().strftime("%Y-%m-%d %I:%M %p"),
            "test_ts": datetime.utcnow().isoformat() + "Z",
        },
    )
    request.session["results"] = results
    return JsonResponse({"ok": True})


@require_GET
def results(request):
    return render(
        request, "results.html", {"results": request.session.get("results", [])}
    )


@require_POST
def password_reset_request(request):
    form = PasswordResetForm(request.POST)
    if form.is_valid():
        form.save(
            request=request,
            use_https=request.is_secure(),
            email_template_name="registration/password_reset_email.txt",
            html_email_template_name="registration/password_reset_email.html",
            subject_template_name="registration/password_reset_subject.txt",
            from_email=settings.DEFAULT_FROM_EMAIL,
        )
    # Always return ok to prevent email enumeration
    return JsonResponse({"status": "ok"})


def password_reset_confirm(request, uidb64, token):
    try:
        uid = force_str(urlsafe_base64_decode(uidb64))
        user = User.objects.get(pk=uid)
    except (User.DoesNotExist, ValueError, TypeError):
        user = None

    valid = user is not None and default_token_generator.check_token(user, token)

    if request.method == "POST" and valid:
        form = SetPasswordForm(user, request.POST)
        if form.is_valid():
            form.save()
            return redirect("password_reset_complete")
        return render(request, "registration/password_reset_confirm.html", {"form": form, "valid": True})

    return render(request, "registration/password_reset_confirm.html", {
        "form": SetPasswordForm(user) if valid else None,
        "valid": valid,
        "uidb64": uidb64,
        "token": token,
    })


def password_reset_complete(request):
    return render(request, "registration/password_reset_complete.html")



@require_POST
def send_feedback(request):
    rating = request.POST.get("rating", "N/A")
    category = request.POST.get("category", "N/A")
    message = (request.POST.get("message", "") or "").strip()

    if not message:
        return JsonResponse({"status": "error", "message": "Message cannot be empty."})

    try:
        subject = f"Studyassists Feedback — {category}"
        body = f"Rating: {rating}\nCategory: {category}\n\nMessage:\n{message}"

        send_mail(
            subject=subject,
            message=body,
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=["info@studyassists.com"],
            fail_silently=False,
        )

        return JsonResponse(
            {"status": "success", "message": "Thank you! Your feedback has been sent."}
        )
    except Exception as e:
        # log the full exception so we can inspect it in the console/logs
        import logging

        logging.getLogger(__name__).exception("feedback send failed")
        # tell the client what went wrong (remove before production)
        return JsonResponse(
            {"status": "error", "message": f"Unable to send feedback: {e}"}
        )



# ---------- Google OAuth2 ----------

GOOGLE_AUTH_URL = "https://accounts.google.com/o/oauth2/v2/auth"
GOOGLE_TOKEN_URL = "https://oauth2.googleapis.com/token"
GOOGLE_USERINFO_URL = "https://www.googleapis.com/oauth2/v3/userinfo"


@never_cache
@require_GET
def google_login(request):
    state = secrets.token_urlsafe(16)
    request.session["google_oauth_state"] = state

    log.info("Google OAuth redirect_uri: %s", settings.GOOGLE_REDIRECT_URI)

    params = {
        "client_id": settings.GOOGLE_CLIENT_ID,
        "redirect_uri": settings.GOOGLE_REDIRECT_URI,
        "response_type": "code",
        "scope": "openid email profile",
        "state": state,
        "access_type": "online",
    }
    return redirect(f"{GOOGLE_AUTH_URL}?{urllib.parse.urlencode(params)}")


@never_cache
@require_GET
def google_callback(request):
    state = request.GET.get("state")
    stored_state = request.session.pop("google_oauth_state", None)

    if not state or state != stored_state:
        messages.error(request, "Authentication failed. Please try again.", extra_tags="auth-login")
        return redirect("/")

    code = request.GET.get("code")
    if not code:
        messages.error(request, "Google sign-in was cancelled.", extra_tags="auth-login")
        return redirect("/")

    # Exchange authorisation code for tokens
    token_resp = http_requests.post(GOOGLE_TOKEN_URL, data={
        "code": code,
        "client_id": settings.GOOGLE_CLIENT_ID,
        "client_secret": settings.GOOGLE_CLIENT_SECRET,
        "redirect_uri": settings.GOOGLE_REDIRECT_URI,
        "grant_type": "authorization_code",
    }, timeout=10)

    if not token_resp.ok:
        messages.error(request, "Could not complete Google sign-in. Please try again.", extra_tags="auth-login")
        return redirect("/")

    access_token = token_resp.json().get("access_token")

    # Fetch user profile from Google
    userinfo_resp = http_requests.get(
        GOOGLE_USERINFO_URL,
        headers={"Authorization": f"Bearer {access_token}"},
        timeout=10,
    )

    if not userinfo_resp.ok:
        messages.error(request, "Could not retrieve profile from Google.", extra_tags="auth-login")
        return redirect("/")

    info = userinfo_resp.json()
    email = (info.get("email") or "").strip().lower()

    if not email:
        messages.error(request, "Google did not return an email address.", extra_tags="auth-login")
        return redirect("/")

    # Get existing user or create a new one
    user = User.objects.filter(email__iexact=email).first()
    if not user:
        first_name = info.get("given_name", "")
        last_name = info.get("family_name", "")
        user = User(
            username=email,
            email=email,
            first_name=first_name,
            last_name=last_name,
        )
        user.set_unusable_password()
        user.save()

    login(request, user)
    return redirect("/")


